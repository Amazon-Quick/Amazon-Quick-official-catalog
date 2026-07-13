"""
Phase 2, Document reconstruction for the translate-ms skill.

Handles translated document XML reconstruction in code rather than by the LLM to ensure deterministic
placement of translated text into the original XML structure. Reduces expensive reasoning
steps and guaranteeing formatting fidelity at any document size.

Provides store_translation_result() which parses worker output (handles malformed
JSON, curly quotes, truncated responses), and reconstruct_document() which applies
translated text back into the original document preserving all formatting.

Also provides check_translation_coverage() and backfill_missing_translations() for
handling failed worker batches.

Requires _STATE from state.py (populated during Phase 1).

File layout:
  scripts/
  ├── state.py         ← TranslationState class + _STATE singleton
  ├── extract.py       ← document extraction (Phase 1)
  ├── glossary.py      ← custom glossary CSV parsing and matching
  └── reconstruct.py   ← this file (reconstruction for both DOCX and PPTX)
"""

import json
import os
import re
import shutil

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation

# try/except avoids F821 lint errors in Quick's script concatenation mode.
try:
    from extract import _get_all_runs_in_order
    from state import _STATE
except ImportError:
    pass  # Names provided by script concatenation at runtime


def _lenient_json_parse(text):
    """2-pass JSON parser for worker responses (clean+parse, then per-object fallback)."""
    if not text or not text.strip():
        return None

    text = text.strip()
    text = re.sub(r"^```(?:json)?\s*\n?", "", text)
    text = re.sub(r"\n?```\s*$", "", text)

    try:
        start = text.index("[")
        end = text.rindex("]") + 1
        text = text[start:end]
    except ValueError:
        return None

    # Normalize smart/curly quotes (including German „ U+201E)
    text = text.replace("\u201c", '"').replace("\u201d", '"').replace("\u201e", '"')
    text = text.replace("\u2018", "'").replace("\u2019", "'")

    # Try json.loads FIRST before any aggressive text mangling
    try:
        result = json.loads(text)
        if isinstance(result, list):
            return result
    except (json.JSONDecodeError, TypeError):
        pass

    # Only apply quote-escaping regex as a recovery step for malformed JSON
    text = re.sub(
        r'("text"\s*:\s*")(.*?)("(?:\s*[,}\]]))',
        lambda m: m.group(1) + m.group(2).replace('"', '\\"') + m.group(3),
        text,
        flags=re.DOTALL,
    )

    try:
        result = json.loads(text)
        if isinstance(result, list):
            return result
    except (json.JSONDecodeError, TypeError):
        pass

    # Fallback: per-object extraction
    try:
        para_blocks = re.findall(
            r'\{\s*"para_id"\s*:\s*(\d+)\s*,\s*"runs"\s*:\s*(\[.*?\])\s*\}',
            text,
            re.DOTALL,
        )
        if para_blocks:
            results = []
            for para_id_str, runs_raw in para_blocks:
                try:
                    runs = json.loads(runs_raw)
                    results.append({"para_id": int(para_id_str), "runs": runs})
                    continue
                except json.JSONDecodeError:
                    pass
                run_items = re.findall(
                    r'\{\s*"id"\s*:\s*(\d+)\s*,\s*"text"\s*:\s*"((?:[^"\\\\]|\\\\.)*)"\s*\}',
                    runs_raw,
                )
                if run_items:
                    runs = [{"id": int(rid), "text": rtxt} for rid, rtxt in run_items]
                    results.append({"para_id": int(para_id_str), "runs": runs})
                else:
                    texts = re.findall(r'"text"\s*:\s*"([^"]*)"', runs_raw)
                    if texts:
                        runs = [{"id": i, "text": t} for i, t in enumerate(texts)]
                        results.append({"para_id": int(para_id_str), "runs": runs})
            if results:
                return results
    except (ValueError, re.error):
        pass

    return None


def _split_text_at_word_boundaries(text, num_parts):
    """Split text into N roughly-equal parts at word boundaries."""
    if num_parts <= 1:
        return [text]
    words = text.split(" ")
    if len(words) <= num_parts:
        parts = [w + " " for w in words[:-1]] + [words[-1]] if words else [""]
        while len(parts) < num_parts:
            parts.append("")
        return parts
    words_per_part = len(words) / num_parts
    parts = []
    start = 0
    for i in range(num_parts):
        end = round(words_per_part * (i + 1))
        chunk_words = words[start:end]
        chunk = " ".join(chunk_words)
        if i < num_parts - 1:
            chunk += " "
        parts.append(chunk)
        start = end
    return parts


def store_translation_result(raw_text):
    """Parse and store a worker result. Handles JSON, code blocks, and error cases."""
    if not raw_text or not raw_text.strip():
        print("WARNING: Empty result received")
        return

    parsed = _lenient_json_parse(raw_text)
    if parsed is None:
        print(
            f"WARNING: Could not parse result ({len(raw_text)} chars). "
            f"First 100: {raw_text[:100]} "
        )
        return

    if not isinstance(parsed, list):
        parsed = [parsed]

    count = 0
    for item in parsed:
        if isinstance(item, dict) and "para_id" in item and "runs" in item:
            pid = item["para_id"]
            runs = item["runs"]
            _STATE.translations[int(pid)] = runs
            count += 1

    _STATE._batch_results_received += 1
    print(
        f"Batch {_STATE._batch_results_received}: stored {count} paragraphs "
        f"(total: {len(_STATE.translations)})"
    )


def check_translation_coverage():
    """Check how many paragraphs were translated and return list of missing para_ids."""
    total = _STATE.stats.get("total_paragraphs", 0)
    translated = len(_STATE.translations)
    missing = []
    for batch in _STATE.batches:
        for p in batch:
            if p["para_id"] not in _STATE.translations:
                missing.append(p["para_id"])
    print(f"Coverage: {translated}/{total}")
    if missing:
        print(f"Missing: {len(missing)} paragraphs, IDs: {missing}")
    else:
        print("All paragraphs translated successfully.")
    return missing


def backfill_missing_translations():
    """Fall back to source text for any paragraphs still not in _STATE.translations."""
    backfilled = 0
    for batch in _STATE.batches:
        for p in batch:
            if p["para_id"] not in _STATE.translations:
                _STATE.translations[p["para_id"]] = p["runs"]
                backfilled += 1
    if backfilled:
        print(f"Backfilled {backfilled} paragraphs with source text")
    return backfilled


def get_missing_paragraph_texts():
    """Return list of (para_id, full_paragraph, runs) for untranslated paragraphs."""
    missing = []
    for batch in _STATE.batches:
        for p in batch:
            if p["para_id"] not in _STATE.translations:
                missing.append((p["para_id"], p["full_paragraph"], p["runs"]))
    return missing


def reconstruct_document(target_language, workspace_dir):
    """PUBLIC. Builds output path, applies translations, saves result."""
    if not _STATE.translations:
        raise ValueError(
            "Cannot reconstruct: zero translations stored. "
            "Ensure workers completed and store_translation_result() was called."
        )

    # Build output path and copy source file (shared across formats)
    ext = os.path.splitext(_STATE.file_path)[1].lower()
    base_name = os.path.splitext(os.path.basename(_STATE.file_path))[0]
    lang_suffix = target_language.lower().strip().replace(" ", "-")
    output_filename = f"{base_name}-{lang_suffix}{ext}"
    output_dir = os.path.join(workspace_dir, "artifacts")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, output_filename)
    shutil.copy2(_STATE.file_path, output_path)

    # Dispatch to format-specific handler
    if _STATE.doc_format == "docx":
        doc, applied, skipped, errors = _reconstruct_docx(output_path)
        doc.save(output_path)
    elif _STATE.doc_format == "pptx":
        prs, applied, skipped, errors = _reconstruct_pptx(output_path)
        prs.save(output_path)
    else:
        raise ValueError(f"No document loaded (doc_format={_STATE.doc_format})")

    print(f"Saved: {output_path}")
    print(f"  Translated: {applied} | Skipped: {skipped} | Errors: {errors}")
    return output_path


# ─────────────────────────────────────────────────────────────────────────────
# DOCX reconstruction
# ─────────────────────────────────────────────────────────────────────────────


def _set_run_text_on_elem(r_elem, new_text):
    """Set text on a run element, preserving xml:space for leading/trailing spaces."""
    t_elems = r_elem.findall(qn("w:t"))
    if t_elems:
        t_elems[0].text = new_text
        if new_text and (new_text[0] == " " or new_text[-1] == " "):
            t_elems[0].set(qn("xml:space"), "preserve")
        for t in t_elems[1:]:
            t.text = ""


def _apply_translations_to_paragraph_docx(paragraph, translated_runs, run_map_entry):
    """Apply translated text to a DOCX paragraph's runs (including hyperlink runs)."""
    original_runs = _get_all_runs_in_order(paragraph)
    merged_run_map = run_map_entry["merged_run_map"]
    original_run_texts = {}
    for run_id, run_data in enumerate(translated_runs):
        text = run_data.get("text", "")
        orig_indices = merged_run_map.get(
            run_id, merged_run_map.get(str(run_id), [run_id])
        )
        num_original = len(orig_indices)
        if num_original == 1:
            original_run_texts[orig_indices[0]] = text
        elif num_original == 0:
            original_run_texts[merged_run_map.get(run_id, [run_id])[0]] = text
        else:
            parts = _split_text_at_word_boundaries(text, num_original)
            for i, orig_idx in enumerate(orig_indices):
                part = parts[i] if i < len(parts) else ""
                original_run_texts[orig_idx] = part
    for idx, r_elem in enumerate(original_runs):
        if idx in original_run_texts:
            _set_run_text_on_elem(r_elem, original_run_texts[idx])


def _reconstruct_docx(output_path):
    """Apply translations to a DOCX file. Returns (Document, applied, skipped, errors)."""
    doc = Document(output_path)
    all_translations = _STATE.translations
    run_map = _STATE.run_map
    applied = skipped = errors = 0
    para_id = 0

    for paragraph in doc.paragraphs:
        if (
            paragraph.text
            and paragraph.text.strip()
            and _get_all_runs_in_order(paragraph)
        ):
            if para_id in run_map and para_id in all_translations:
                try:
                    _apply_translations_to_paragraph_docx(
                        paragraph, all_translations[para_id], run_map[para_id]
                    )
                    applied += 1
                except Exception as e:
                    errors += 1
                    print(f"Error applying para {para_id}: {e}")
            else:
                skipped += 1
            para_id += 1
        elif paragraph.text and paragraph.text.strip():
            para_id += 1

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if (
                        paragraph.text
                        and paragraph.text.strip()
                        and _get_all_runs_in_order(paragraph)
                    ):
                        if para_id in run_map and para_id in all_translations:
                            try:
                                _apply_translations_to_paragraph_docx(
                                    paragraph,
                                    all_translations[para_id],
                                    run_map[para_id],
                                )
                                applied += 1
                            except Exception as e:
                                errors += 1
                                print(f"Error applying para {para_id}: {e}")
                        else:
                            skipped += 1
                        para_id += 1
                    elif paragraph.text and paragraph.text.strip():
                        para_id += 1

    for section in doc.sections:
        for header_footer in [
            section.header,
            section.footer,
            section.first_page_header,
            section.first_page_footer,
        ]:
            if header_footer and header_footer.is_linked_to_previous is False:
                for paragraph in header_footer.paragraphs:
                    if (
                        paragraph.text
                        and paragraph.text.strip()
                        and _get_all_runs_in_order(paragraph)
                    ):
                        if para_id in run_map and para_id in all_translations:
                            try:
                                _apply_translations_to_paragraph_docx(
                                    paragraph,
                                    all_translations[para_id],
                                    run_map[para_id],
                                )
                                applied += 1
                            except Exception as e:
                                errors += 1
                                print(f"Error applying para {para_id}: {e}")
                        else:
                            skipped += 1
                        para_id += 1
                    elif paragraph.text and paragraph.text.strip():
                        para_id += 1

    return doc, applied, skipped, errors


# ─────────────────────────────────────────────────────────────────────────────
# PPTX reconstruction
# ─────────────────────────────────────────────────────────────────────────────


def _apply_translations_to_paragraph_pptx(paragraph, translated_runs, run_map_entry):
    """Apply translated text to a PPTX paragraph's runs."""
    original_runs = paragraph.runs
    merged_run_map = run_map_entry["merged_run_map"]
    original_run_texts = {}
    for run_id, run_data in enumerate(translated_runs):
        text = run_data.get("text", "")
        orig_indices = merged_run_map.get(
            run_id, merged_run_map.get(str(run_id), [run_id])
        )
        num_original = len(orig_indices)
        if num_original == 1:
            original_run_texts[orig_indices[0]] = text
        else:
            parts = _split_text_at_word_boundaries(text, num_original)
            for i, orig_idx in enumerate(orig_indices):
                part = parts[i] if i < len(parts) else ""
                original_run_texts[orig_idx] = part
    for idx, run in enumerate(original_runs):
        if idx in original_run_texts:
            run.text = original_run_texts[idx]


def _reconstruct_text_frame(text_frame, para_id):
    """Apply translations to all paragraphs in a text frame."""
    all_translations = _STATE.translations
    run_map = _STATE.run_map
    applied = skipped = errors = 0
    for paragraph in text_frame.paragraphs:
        if not paragraph.text.strip():
            para_id += 1
            continue
        if para_id in run_map and para_id in all_translations:
            try:
                _apply_translations_to_paragraph_pptx(
                    paragraph, all_translations[para_id], run_map[para_id]
                )
                applied += 1
            except Exception as e:
                errors += 1
                print(f"Error applying para {para_id}: {e}")
        else:
            skipped += 1
        para_id += 1
    return para_id, applied, skipped, errors


def _iter_pptx_shapes(shapes):
    """Recursively yield all shapes, walking into GroupShape containers."""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child in shape.shapes:
                    yield from _iter_pptx_shapes([child])
            except AttributeError:
                pass
        else:
            yield shape


def _reconstruct_pptx(output_path):
    """Apply translations to a PPTX file. Returns (Presentation, applied, skipped, errors)."""
    prs = Presentation(output_path)
    applied = skipped = errors = 0
    para_id = 0

    for slide in prs.slides:
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.has_text_frame:
                para_id, a, s, e = _reconstruct_text_frame(shape.text_frame, para_id)
                applied += a
                skipped += s
                errors += e
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            para_id, a, s, e = _reconstruct_text_frame(
                                cell.text_frame, para_id
                            )
                            applied += a
                            skipped += s
                            errors += e
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            para_id, a, s, e = _reconstruct_text_frame(
                slide.notes_slide.notes_text_frame, para_id
            )
            applied += a
            skipped += s
            errors += e

    return prs, applied, skipped, errors
