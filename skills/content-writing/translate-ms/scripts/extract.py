"""
Phase 1, Document extraction for the translate-ms skill.

Offloads XML parsing, run merging, and state management to deterministic code
(rather than LLM reasoning) eliminates expensive inference steps, enables reliable
handling of large documents, and produces consistent results regardless of model context limits.

Provides extract_document() which detects .docx/.pptx, walks the document structure
(body, tables, headers/footers, shapes, notes), merges runs with identical formatting,
and batches paragraphs for parallel translation by workers. Also provides
get_first_words() for language detection.

Loaded in Phase 1 before spawning workers. Receives _STATE from state.py.

File layout:
  scripts/
  ├── state.py         ← TranslationState class + _STATE singleton
  ├── extract.py       ← this file (extraction for both DOCX and PPTX)
  ├── glossary.py      ← custom glossary CSV parsing and matching
  └── reconstruct.py   ← worker result parsing + document reconstruction
"""

import os

from docx import Document
from docx.oxml.ns import qn
from lxml import etree
from pptx import Presentation

# try/except avoids F821 lint errors in Quick's script concatenation mode.
try:
    from state import _STATE
except ImportError:
    pass  # Name provided by script concatenation at runtime


# ─────────────────────────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────────────────────────

WORDS_PER_BATCH = 250
MAX_PARAS_PER_BATCH = 30


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────


def get_first_words(n=100):
    """Return the first n words from batch data (for language detection)."""
    words = []
    for batch in _STATE.batches:
        for p in batch:
            words.extend(p["full_paragraph"].split())
            if len(words) >= n:
                return " ".join(words[:n])
    return " ".join(words)


def extract_document(file_path):
    """PUBLIC. Detect format, extract, batch. Returns (batches, stats)."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".docx":
        _STATE.reset("docx", file_path)
        all_paragraphs = _extract_docx(file_path)
    elif ext == ".pptx":
        _STATE.reset("pptx", file_path)
        all_paragraphs = _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

    _STATE.batches = _batch_paragraphs(all_paragraphs)
    batch_count = len(_STATE.batches)
    _STATE.stats["batch_count"] = batch_count

    print(f"Extracted: {_STATE.stats['total_paragraphs']} paragraphs")
    stat_parts = [
        f"{k}: {v}"
        for k, v in _STATE.stats.items()
        if k != "total_paragraphs" and k != "batch_count"
    ]
    print(f"  {' | '.join(stat_parts)}")
    print(f"  Batches: {batch_count}")

    return _STATE.batches, _STATE.stats


# ─────────────────────────────────────────────────────────────────────────────
# Shared helpers
# ─────────────────────────────────────────────────────────────────────────────


def _batch_paragraphs(all_paragraphs):
    """Group paragraphs into batches by word count."""
    batches = []
    current_batch = []
    current_words = 0
    for p in all_paragraphs:
        p_words = len(p["full_paragraph"].split())
        if current_batch and (
            current_words + p_words > WORDS_PER_BATCH
            or len(current_batch) >= MAX_PARAS_PER_BATCH
        ):
            batches.append(current_batch)
            current_batch = []
            current_words = 0
        current_batch.append(p)
        current_words += p_words
    if current_batch:
        batches.append(current_batch)
    return batches


# ─────────────────────────────────────────────────────────────────────────────
# DOCX extraction
# ─────────────────────────────────────────────────────────────────────────────


def _get_all_runs_in_order(paragraph):
    """Get ALL run elements including those inside <w:hyperlink> elements."""
    runs = []
    for child in paragraph._p:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag == "r":
            runs.append(child)
        elif tag == "hyperlink":
            for sub in child:
                sub_tag = (
                    etree.QName(sub.tag).localname if isinstance(sub.tag, str) else ""
                )
                if sub_tag == "r":
                    runs.append(sub)
    return runs


def _get_run_text(r_elem):
    """Get text from a run element."""
    return "".join(t.text or "" for t in r_elem.findall(qn("w:t")))


def _get_formatting_key(r_elem):
    """Get a hashable key representing a run's formatting (canonical rPr XML)."""
    rpr = r_elem.find(qn("w:rPr"))
    if rpr is None:
        return ""
    return etree.tostring(rpr, method="c14n").decode()


def _merge_docx_runs(runs):
    """Merge adjacent runs with identical formatting. Returns [(text, [indices])]."""
    if not runs:
        return []
    merged = []
    current_text = _get_run_text(runs[0])
    current_key = _get_formatting_key(runs[0])
    current_indices = [0]
    for i in range(1, len(runs)):
        key = _get_formatting_key(runs[i])
        if key == current_key:
            current_text += _get_run_text(runs[i])
            current_indices.append(i)
        else:
            merged.append((current_text, current_indices[:]))
            current_text = _get_run_text(runs[i])
            current_key = key
            current_indices = [i]
    merged.append((current_text, current_indices[:]))
    return merged


def _process_docx_paragraph(paragraph, para_id, all_paragraphs):
    """Process a single DOCX paragraph into extraction format."""
    runs = _get_all_runs_in_order(paragraph)
    if not runs:
        return para_id
    merged = _merge_docx_runs(runs)
    full_text = "".join(text for text, _ in merged)
    if not full_text.strip():
        return para_id
    run_list = []
    merged_run_map = {}
    for run_id, (text, orig_indices) in enumerate(merged):
        run_list.append({"id": run_id, "text": text})
        merged_run_map[run_id] = orig_indices
    _STATE.run_map[para_id] = {
        "merged_run_map": merged_run_map,
        "run_count": len(runs),
    }
    all_paragraphs.append(
        {"para_id": para_id, "full_paragraph": full_text, "runs": run_list}
    )
    return para_id + 1


def _extract_docx(file_path):
    """Walk DOCX body, tables, headers/footers. Returns all_paragraphs."""
    doc = Document(file_path)
    all_paragraphs = []
    para_id = 0

    for paragraph in doc.paragraphs:
        if (
            paragraph.text
            and paragraph.text.strip()
            and _get_all_runs_in_order(paragraph)
        ):
            para_id = _process_docx_paragraph(paragraph, para_id, all_paragraphs)
        elif paragraph.text and paragraph.text.strip():
            para_id += 1

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if (
                        paragraph.text
                        and paragraph.text.strip()
                        and _get_all_runs_in_order(paragraph)
                    ):
                        para_id = _process_docx_paragraph(
                            paragraph, para_id, all_paragraphs
                        )
                    elif paragraph.text and paragraph.text.strip():
                        para_id += 1

    for section in doc.sections:
        for header_footer in [
            section.header,
            section.footer,
            section.first_page_header,
            section.first_page_footer,
        ]:
            if header_footer and header_footer.is_linked_to_previous is False:
                for paragraph in header_footer.paragraphs:
                    if (
                        paragraph.text
                        and paragraph.text.strip()
                        and _get_all_runs_in_order(paragraph)
                    ):
                        para_id = _process_docx_paragraph(
                            paragraph, para_id, all_paragraphs
                        )
                    elif paragraph.text and paragraph.text.strip():
                        para_id += 1

    _STATE.stats = {
        "total_paragraphs": len(all_paragraphs),
    }
    return all_paragraphs


# ─────────────────────────────────────────────────────────────────────────────
# PPTX extraction
# ─────────────────────────────────────────────────────────────────────────────


def _get_pptx_font_key(run):
    """Get a hashable key representing a PPTX run's formatting."""
    font = run.font
    try:
        color_val = str(font.color.rgb) if font.color and font.color.rgb else None
    except (AttributeError, TypeError):
        color_val = None
    return (
        str(font.name),
        str(font.size),
        str(font.bold),
        str(font.italic),
        str(font.underline),
        str(color_val),
    )


def _merge_pptx_runs(runs):
    """Merge adjacent PPTX runs with identical font properties."""
    if not runs:
        return []
    merged = []
    current_text = runs[0].text or ""
    current_font = _get_pptx_font_key(runs[0])
    current_indices = [0]
    for i in range(1, len(runs)):
        font_key = _get_pptx_font_key(runs[i])
        if font_key == current_font:
            current_text += runs[i].text or ""
            current_indices.append(i)
        else:
            merged.append((current_text, current_indices[:]))
            current_text = runs[i].text or ""
            current_font = font_key
            current_indices = [i]
    merged.append((current_text, current_indices[:]))
    return merged


def _process_pptx_text_frame(text_frame, para_id, all_paragraphs):
    """Process all paragraphs in a PPTX text frame."""
    for paragraph in text_frame.paragraphs:
        if not paragraph.text.strip():
            para_id += 1
            continue
        runs = paragraph.runs
        if not runs:
            all_paragraphs.append(
                {
                    "para_id": para_id,
                    "full_paragraph": paragraph.text,
                    "runs": [{"id": 0, "text": paragraph.text}],
                }
            )
            _STATE.run_map[para_id] = {"merged_run_map": {0: [0]}, "run_count": 1}
            para_id += 1
            continue
        merged = _merge_pptx_runs(runs)
        full_text = "".join(text for text, _ in merged)
        if not full_text.strip():
            para_id += 1
            continue
        run_list = []
        merged_run_map = {}
        for run_id, (text, orig_indices) in enumerate(merged):
            run_list.append({"id": run_id, "text": text})
            merged_run_map[run_id] = orig_indices
        _STATE.run_map[para_id] = {
            "merged_run_map": merged_run_map,
            "run_count": len(runs),
        }
        all_paragraphs.append(
            {"para_id": para_id, "full_paragraph": full_text, "runs": run_list}
        )
        para_id += 1
    return para_id


def _iter_pptx_shapes(shapes):
    """Recursively yield all shapes, walking into GroupShape containers."""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child in shape.shapes:
                    yield from _iter_pptx_shapes([child])
            except AttributeError:
                pass
        else:
            yield shape


def _extract_pptx(file_path):
    """Walk PPTX slides, shapes, tables, notes. Returns all_paragraphs."""
    prs = Presentation(file_path)
    all_paragraphs = []
    para_id = 0
    slide_count = 0

    for slide in prs.slides:
        slide_count += 1
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.has_text_frame:
                para_id = _process_pptx_text_frame(
                    shape.text_frame, para_id, all_paragraphs
                )
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            para_id = _process_pptx_text_frame(
                                cell.text_frame, para_id, all_paragraphs
                            )
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            para_id = _process_pptx_text_frame(
                slide.notes_slide.notes_text_frame, para_id, all_paragraphs
            )

    _STATE.stats = {
        "total_paragraphs": len(all_paragraphs),
        "slides": slide_count,
    }
    return all_paragraphs
